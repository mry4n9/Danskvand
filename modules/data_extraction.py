import streamlit as st
import requests
from bs4 import BeautifulSoup
from pypdf import PdfReader
from pptx import Presentation
import io

@st.cache_data(show_spinner=False)
def extract_text_from_url(_url: str) -> str | None:
    """Extracts all text content from a URL."""
    try:
        response = requests.get(_url, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove script and style elements
        for script_or_style in soup(["script", "style"]):
            script_or_style.decompose()
        
        text = soup.get_text(separator=' ', strip=True)
        return text
    except requests.exceptions.RequestException as e:
        st.error(f"Error fetching URL {_url}: {e}")
        return None
    except Exception as e:
        st.error(f"Error parsing URL content: {e}")
        return None

@st.cache_data(show_spinner=False)
def extract_text_from_pdf_bytes(_pdf_bytes: bytes) -> str | None:
    """Extracts text from PDF bytes."""
    try:
        reader = PdfReader(io.BytesIO(_pdf_bytes))
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        st.error(f"Error reading PDF file: {e}")
        return None

@st.cache_data(show_spinner=False)
def extract_text_from_pptx_bytes(_pptx_bytes: bytes) -> str | None:
    """Extracts text from PPTX bytes."""
    try:
        prs = Presentation(io.BytesIO(_pptx_bytes))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"Error reading PPTX file: {e}")
        return None

def extract_text_from_file(uploaded_file) -> str | None:
    """Detects file type and extracts text."""
    if uploaded_file is None:
        return None
    
    file_bytes = uploaded_file.getvalue()
    if uploaded_file.type == "application/pdf":
        return extract_text_from_pdf_bytes(file_bytes)
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_pptx_bytes(file_bytes)
    else:
        st.warning(f"Unsupported file type: {uploaded_file.type}")
        return None